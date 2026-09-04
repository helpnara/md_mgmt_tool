# -*- coding: utf-8 -*-
"""사내 AX 성과공유회 제출용 세로형 포스터(A1)를 만든다.

    python make_poster.py            # docs/poster/ 에 pptx 생성
    soffice --headless --convert-to pdf ...   # PDF 는 별도 변환

숫자는 전부 2026-09-04 실측치이며, 근거는 docs/외주개발-비용산정.md 에 있다.
수치를 고칠 일이 생기면 아래 DATA 사전만 고치면 된다.
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Mm, Pt

OUT = Path(__file__).resolve().parent

# ── 색 ─────────────────────────────────────────────────────────────────────
NAVY   = RGBColor(0x0F, 0x2A, 0x47)
NAVY_D = RGBColor(0x09, 0x1B, 0x2E)
SKY    = RGBColor(0x9C, 0xC3, 0xEE)
BLUE   = RGBColor(0x1F, 0x6F, 0xEB)
BLUE_L = RGBColor(0xEB, 0xF2, 0xFE)
TEAL   = RGBColor(0x0B, 0x8A, 0x5F)
TEAL_L = RGBColor(0xE4, 0xF5, 0xEE)
AMBER  = RGBColor(0xA1, 0x59, 0x0B)
AMBER_L= RGBColor(0xFD, 0xF3, 0xE3)
INK    = RGBColor(0x0F, 0x17, 0x2A)
MUTED  = RGBColor(0x4B, 0x5A, 0x6E)
FAINT  = RGBColor(0x7A, 0x88, 0x99)
LINE   = RGBColor(0xC9, 0xD4, 0xE0)
SURF   = RGBColor(0xF7, 0xF9, 0xFB)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)

FONT = "맑은 고딕"

# ── 실측 수치 (2026-09-04) ────────────────────────────────────────────────
DATA = dict(
    days="3일", hours="30시간", mm_real="0.19 M/M", cost_real="약 250만 원",
    mm_out="16 M/M", period_out="4~5개월", cost_out="약 2.0억 원",
    saving="약 1억 9,750만 원", ratio="약 80배", fp="319 FP",
    loc="12,876줄", api="56개", tests="382건", screens="6개",
    tables="9개", todo_all="69건", todo_done="59건", todo_wait="10건",
)

# ── 조판 상수 (mm) ────────────────────────────────────────────────────────
PW, PH = 594.0, 841.0          # A1 세로
M = 30.0                        # 좌우 여백
W = PW - 2 * M                  # 본문 폭


# ── 그리기 도우미 ─────────────────────────────────────────────────────────
def rect(slide, x, y, w, h, fill=None, line=None, radius=None, lw=0.8):
    shape_kind = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    s = slide.shapes.add_shape(shape_kind, Mm(x), Mm(y), Mm(w), Mm(h))
    if radius:
        s.adjustments[0] = radius
    if fill is None:
        s.fill.background()
    else:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = Pt(lw)
    s.shadow.inherit = False
    s.text_frame.word_wrap = True
    return s


def text(slide, x, y, w, h, lines, anchor=MSO_ANCHOR.TOP, align=PP_ALIGN.LEFT):
    """lines: [{t, size, bold, color, align, before, after, spacing}] 목록."""
    box = slide.shapes.add_textbox(Mm(x), Mm(y), Mm(w), Mm(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, spec in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = spec.get("align", align)
        p.space_before = Pt(spec.get("before", 0))
        p.space_after = Pt(spec.get("after", 0))
        if spec.get("spacing"):
            p.line_spacing = spec["spacing"]
        run = p.add_run()
        run.text = spec["t"]
        f = run.font
        f.name = FONT
        f.size = Pt(spec.get("size", 20))
        f.bold = spec.get("bold", False)
        f.color.rgb = spec.get("color", INK)
    return box


def section_title(slide, y, number, title, note=""):
    """섹션 머리. 왼쪽에 번호 원, 오른쪽에 제목과 보조 설명."""
    d = 13.0
    c = rect(slide, M, y, d, d, fill=NAVY, radius=0.5)
    text(slide, M, y + 2.3, d, d, [{"t": number, "size": 20, "bold": True,
                                    "color": WHITE, "align": PP_ALIGN.CENTER}])
    text(slide, M + d + 6, y - 0.6, W - d - 6, d + 4,
         [{"t": title, "size": 30, "bold": True, "color": NAVY}])
    if note:
        text(slide, M, y + d + 3.2, W, 8,
             [{"t": note, "size": 16, "color": MUTED}])
    return y + d + (13 if note else 6)


def bullets(slide, x, y, w, items, size=16.5, color=INK, gap=3.0, mark="·"):
    lines = []
    for i, it in enumerate(items):
        lines.append({"t": f"{mark} {it}", "size": size, "color": color,
                      "after": gap if i < len(items) - 1 else 0, "spacing": 1.18})
    text(slide, x, y, w, 200, lines)


# ── 포스터 ────────────────────────────────────────────────────────────────
def build() -> Path:
    prs = Presentation()
    prs.slide_width, prs.slide_height = Mm(PW), Mm(PH)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    rect(slide, 0, 0, PW, PH, fill=WHITE)

    # ── 머리말 ────────────────────────────────────────────────────────────
    HDR = 122.0
    rect(slide, 0, 0, PW, HDR, fill=NAVY)
    rect(slide, 0, HDR, PW, 2.2, fill=BLUE)

    chip_w = 118.0
    rect(slide, M, 15, chip_w, 11, fill=NAVY_D, radius=0.5)
    text(slide, M, 17.6, chip_w, 8,
         [{"t": "2026 사내 AX 성과공유회", "size": 15, "bold": True,
           "color": SKY, "align": PP_ALIGN.CENTER}])

    text(slide, M, 32, W, 40,
         [{"t": "개발을 기다리지 않았다", "size": 82, "bold": True, "color": WHITE}])
    text(slide, M, 80, W, 26,
         [{"t": "AI와 함께 만든 현업 업무혁신 — 개발 리소스 없이 3일 만에 구축한",
           "size": 24, "color": SKY, "after": 1.5},
          {"t": "팀 과제관리 · 보고지원 시스템", "size": 24, "bold": True, "color": WHITE}])
    text(slide, M, 106, W, 10,
         [{"t": "선강DX개발팀   |   발표자 권경락   |   2026. 09", "size": 15, "color": SKY}])
    tx0 = M + W * 0.635
    rect(slide, tx0 - 9, 36, 0.8, 62, fill=BLUE)
    text(slide, tx0, 36, W - (tx0 - M), 64, [
        {"t": "만들 사람이", "size": 21, "color": SKY, "after": 2},
        {"t": "매일 쓸 사람이면", "size": 21, "bold": True, "color": WHITE, "after": 2},
        {"t": "요구정의도, 검수도,", "size": 21, "color": SKY, "after": 2},
        {"t": "재작업도 없습니다", "size": 21, "bold": True, "color": WHITE},
    ])

    y = HDR + 12

    # ── 핵심 성과 4칸 ─────────────────────────────────────────────────────
    KH = 78.0
    cards = [
        ("구축 기간", DATA["days"], f"실작업 {DATA['hours']} · 1명", f"외주 추정 {DATA['period_out']}", False),
        ("투입 공수", DATA["mm_real"], "요구·개발·검수 1인", f"외주 추정 {DATA['mm_out']}", False),
        ("소요 비용", DATA["cost_real"], "인건비 + AI 요금", f"외주 추정 {DATA['cost_out']}", False),
        ("절감 효과", DATA["ratio"], DATA["saving"] + " 절감", "비용 기준 98.7%↓", True),
    ]
    gap = 7.0
    cw = (W - gap * 3) / 4
    for i, (label, value, sub, cmp_, hi) in enumerate(cards):
        x = M + i * (cw + gap)
        rect(slide, x, y, cw, KH,
             fill=TEAL_L if hi else SURF,
             line=TEAL if hi else LINE, radius=0.10, lw=1.6 if hi else 1.0)
        rect(slide, x, y, cw, 3.0, fill=TEAL if hi else BLUE)
        text(slide, x + 7, y + 10, cw - 14, 8,
             [{"t": label, "size": 15, "bold": True, "color": TEAL if hi else MUTED}])
        text(slide, x + 7, y + 20, cw - 14, 22,
             [{"t": value, "size": 40 if len(value) > 7 else 46, "bold": True,
               "color": TEAL if hi else NAVY}])
        text(slide, x + 7, y + 46, cw - 14, 10,
             [{"t": sub, "size": 14, "color": INK}])
        rect(slide, x + 7, y + 58.5, cw - 14, 0.5, fill=LINE)
        text(slide, x + 7, y + 62, cw - 14, 10,
             [{"t": cmp_, "size": 14, "bold": True, "color": TEAL if hi else FAINT}])
    y += KH + 12

    # ── 1. 왜 만들었나 ────────────────────────────────────────────────────
    y = section_title(slide, y, "1", "왜 만들었나",
                      "과제 정보는 흩어져 있었고, 무엇을 보고할지는 기억에 의존했다")
    BH = 72.0
    hw = (W - 10) / 2
    for idx, (title, tint, edge, items, tag) in enumerate([
        ("BEFORE", AMBER_L, AMBER, [
            "과제 정보가 개인 메모·메일·엑셀에 흩어져 이력과 자료를 한 번에 볼 수 없음",
            "매주 어떤 과제를 보고할지 판단 기준이 없어 담당자의 기억과 경험에 의존",
            "무엇을 보고했는지 기록이 남지 않아 중복 작성·누락 가능성",
        ], "지금까지"),
        ("AFTER", BLUE_L, BLUE, [
            "과제 1건 = 폴더 1개. 개요·진행이력·첨부·보고가 한자리에 모임",
            "미보고 경과일과 새 진행분으로 이번 주 보고 대상을 자동 추천",
            "확정 보고는 그 시점 스냅샷으로 고정 — 무엇을 보고했는지 그대로 남음",
        ], "지금은"),
    ]):
        x = M + idx * (hw + 10)
        rect(slide, x, y, hw, BH, fill=tint, line=edge, radius=0.09, lw=1.2)
        rect(slide, x, y, 3.2, BH, fill=edge)
        text(slide, x + 12, y + 8, hw - 20, 10,
             [{"t": f"{title}   {tag}", "size": 18, "bold": True, "color": edge}])
        bullets(slide, x + 12, y + 22, hw - 22, items, size=15.5, color=INK, gap=4.5)
    ar = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Mm(M + hw - 4), Mm(y + BH / 2 - 5),
                                Mm(18), Mm(10))
    ar.fill.solid(); ar.fill.fore_color.rgb = NAVY
    ar.line.color.rgb = WHITE; ar.line.width = Pt(2)
    ar.shadow.inherit = False
    y += BH + 11

    # ── 2. 무엇을 만들었나 ────────────────────────────────────────────────
    y = section_title(slide, y, "2", "무엇을 만들었나",
                      "과제 관리부터 주간보고 작성까지, 한 도구 안에서")
    FH = 70.0
    feats = [
        ("①", "과제 중심 통합관리",
         ["개요·날짜별 진행이력·첨부를 한 구조로",
          "원본은 Markdown 파일·폴더 — 특정 프로그램에 종속 안 됨",
          "태그·담당자·상태로 정렬·필터"]),
        ("②", "데이터 기반 보고대상 추천",
         ["마지막 보고 이후 경과일 + 새 진행분으로 필요도 산정",
          "이번 주 확인할 과제를 순서대로 제시",
          "기억이 아니라 기록으로 판단"]),
        ("③", "보고 초안 자동생성 · 이력 고정",
         ["지난 보고 이후 진행분만 모아 초안 생성",
          "확정본은 스냅샷으로 잠금 — 당시 내용 그대로 보존",
          "지난 보고 대비 변경분·보고 이력 검색"]),
        ("④", "사내망 전용 · 오프라인 배포",
         ["데이터가 내 PC 밖으로 나가지 않음 (외부 전송 없음)",
          "인터넷 차단 PC 에서 ZIP 하나로 설치",
          "자동 백업 · 버전 보관 · 휴지통"]),
    ]
    fw = (W - 7 * 3) / 4
    for i, (num, title, items) in enumerate(feats):
        x = M + i * (fw + 7)
        rect(slide, x, y, fw, FH, fill=WHITE, line=LINE, radius=0.08, lw=1.0)
        rect(slide, x, y, fw, 23, fill=SURF, radius=0.0)
        rect(slide, x, y + 23, fw, 0.5, fill=LINE)
        text(slide, x + 7, y + 4, 12, 9,
             [{"t": num, "size": 17, "bold": True, "color": BLUE}])
        text(slide, x + 7, y + 11.5, fw - 14, 11,
             [{"t": title, "size": 16, "bold": True, "color": NAVY, "spacing": 1.1}])
        bullets(slide, x + 7, y + 28, fw - 13, items, size=13.5, color=MUTED, gap=3.2)
    y += FH + 7

    # ── 실제 화면 ────────────────────────────────────────────────────────
    shots = [
        ("screen-list.png", "과제 목록 — 상태 · 담당자 · 마감 · 효과 · 미보고를 한 화면에"),
        ("screen-reports.png", "보고 대상 후보 — 미보고 경과일 순으로 자동 추천"),
        ("screen-detail.png", "과제 상세 — 개요와 진행이력, 지난 보고 지점 표시"),
    ]
    iw, ih, ig = 158.0, 79.0, 7.0
    x0 = M + (W - (iw * 3 + ig * 2)) / 2
    for i, (fname, cap) in enumerate(shots):
        x = x0 + i * (iw + ig)
        path = OUT / fname
        if path.exists():
            slide.shapes.add_picture(str(path), Mm(x), Mm(y), Mm(iw), Mm(ih))
        rect(slide, x, y, iw, ih, fill=None, line=LINE, radius=None, lw=1.0)
        text(slide, x, y + ih + 2.2, iw, 8,
             [{"t": cap, "size": 11.5, "color": FAINT, "align": PP_ALIGN.CENTER}])
    y += ih + 8 + 6

    # ── 실측 규모 띠 ──────────────────────────────────────────────────────
    SH = 25.0
    rect(slide, M, y, W, SH, fill=NAVY, radius=0.12)
    metrics = [("코드", DATA["loc"]), ("API", DATA["api"]), ("자동 시험", DATA["tests"]),
               ("화면", DATA["screens"]), ("저장 테이블", DATA["tables"]),
               ("기능점수", DATA["fp"]),
               ("개선 반영", f"{DATA['todo_done']}/{DATA['todo_all']}")]
    mw = W / len(metrics)
    for i, (k, v) in enumerate(metrics):
        x = M + i * mw
        if i:
            rect(slide, x, y + 6, 0.4, SH - 12, fill=RGBColor(0x2E, 0x4B, 0x6B))
        text(slide, x, y + 4, mw, 8,
             [{"t": k, "size": 12.5, "color": SKY, "align": PP_ALIGN.CENTER}])
        text(slide, x, y + 11.5, mw, 11,
             [{"t": v, "size": 19, "bold": True, "color": WHITE, "align": PP_ALIGN.CENTER}])
    y += SH + 12

    # ── 3. 얼마나 아꼈나 ──────────────────────────────────────────────────
    y = section_title(slide, y, "3", "얼마나 아꼈나 — 외주 개발 대비",
                      "같은 시스템을 외부 인원에게 맡겼을 때의 비용을 기능점수로 산정해 비교")
    EH = 108.0
    rect(slide, M, y, W, EH, fill=SURF, line=LINE, radius=0.06, lw=1.0)

    # 왼쪽 : 비교 표
    tx, tw = M + 12, W * 0.46
    rows = [("구분", "외주 개발 (추정)", "실제 (AI 협업)"),
            ("투입 공수", DATA["mm_out"], DATA["mm_real"]),
            ("개발 기간", f"{DATA['period_out']} · 3~4명", f"{DATA['days']} · 1명"),
            ("소요 비용", DATA["cost_out"], DATA["cost_real"])]
    rh = 14.5
    c1, c2 = tw * 0.30, tw * 0.36
    for r, (a, b, c) in enumerate(rows):
        ry = y + 10 + r * rh
        head = r == 0
        if head:
            rect(slide, tx, ry, tw, rh, fill=NAVY, radius=0.0)
        else:
            rect(slide, tx, ry + rh - 0.4, tw, 0.4, fill=LINE)
        col = WHITE if head else INK
        text(slide, tx + 5, ry + 4.2, c1, 10,
             [{"t": a, "size": 15, "bold": True, "color": WHITE if head else MUTED}])
        text(slide, tx + c1, ry + 4.2, c2, 10,
             [{"t": b, "size": 16, "bold": head, "color": col if head else MUTED,
               "align": PP_ALIGN.CENTER}])
        text(slide, tx + c1 + c2, ry + 4.0, tw - c1 - c2, 10,
             [{"t": c, "size": 17 if not head else 16, "bold": True,
               "color": WHITE if head else TEAL, "align": PP_ALIGN.CENTER}])

    # 오른쪽 : 비용 막대
    bx = M + W * 0.52
    bw = W * 0.42
    text(slide, bx, y + 10, bw, 8,
         [{"t": "개발 비용 비교", "size": 15, "bold": True, "color": MUTED}])
    rect(slide, bx, y + 20, bw, 13, fill=NAVY, radius=0.0)
    text(slide, bx + 6, y + 23.2, bw - 12, 10,
         [{"t": f"외주 개발  {DATA['cost_out']}", "size": 16, "bold": True, "color": WHITE}])
    small = max(bw * 0.0125, 4.0)
    rect(slide, bx, y + 37, bw, 13, fill=RGBColor(0xE6, 0xEA, 0xEF), radius=0.0)
    rect(slide, bx, y + 37, small, 13, fill=TEAL, radius=0.0)
    text(slide, bx + small + 5, y + 40.2, bw - small - 8, 10,
         [{"t": f"실제  {DATA['cost_real']}  (1.3%)", "size": 16, "bold": True, "color": TEAL}])
    rect(slide, bx, y + 53, bw, 19, fill=TEAL_L, line=TEAL, radius=0.12, lw=1.2)
    text(slide, bx, y + 57, bw, 14,
         [{"t": f"{DATA['saving']} 절감  ·  {DATA['ratio']}", "size": 20, "bold": True,
           "color": TEAL, "align": PP_ALIGN.CENTER}])

    # 아래 : 근거와 해석
    rect(slide, M + 12, y + 75, W - 24, 0.5, fill=LINE)
    text(slide, M + 12, y + 79, W - 24, 27, [
        {"t": "산출 근거   CRUD 행렬(엔티티 11 × API 56) → 기능점수 간이법 "
              f"{DATA['fp']} → 기능점수법·코드라인법·상향식 3중 교차검증 = {DATA['mm_out']} "
              "(범위 13~20). 비용은 시장 실거래 단가·대가산정 가이드·FP 단가 세 방식(1.6~2.6억)의 중간값. "
              "실제 비용의 AI 요금은 대화 2,060회·입력 8.3억 토큰 실측 기준 약 100만 원.",
         "size": 13.5, "color": MUTED, "after": 4, "spacing": 1.2},
        {"t": "이 차이의 절반은 AI의 구현 속도이고, 나머지 절반은 요구자·개발자·검수자가 "
              "같은 사람일 때 사라지는 회의·문서·검수 비용입니다 (외주 16 M/M 중 4.7 M/M, 29%).",
         "size": 15, "bold": True, "color": NAVY, "spacing": 1.2},
    ])
    y += EH + 11

    # ── 4. 어떻게 만들었나 + 앞으로 ──────────────────────────────────────
    y = section_title(slide, y, "4", "어떻게 만들었나, 그리고 앞으로")
    LH = 70.0
    lw_ = (W - 10) * 0.52
    rw_ = W - 10 - lw_

    # 왼쪽 : 6단계
    rect(slide, M, y, lw_, LH, fill=WHITE, line=LINE, radius=0.08, lw=1.0)
    text(slide, M + 10, y + 8, lw_ - 20, 8,
         [{"t": "AI 협업 개발 6단계 — 다른 업무에도 그대로 쓸 수 있는 방법론",
           "size": 15, "bold": True, "color": NAVY}])
    steps = ["문제 정의", "업무규칙 설계", "AI 협업 구현", "자동 시험", "실사용", "개선 반영"]
    sgap, scount = 2.6, len(steps)
    sw = (lw_ - 20 - sgap * (scount - 1)) / scount
    for i, s in enumerate(steps):
        sx = M + 10 + i * (sw + sgap)
        rect(slide, sx, y + 20, sw, 13.5, fill=BLUE_L, line=BLUE, radius=0.25, lw=0.9)
        text(slide, sx, y + 23.8, sw, 9,
             [{"t": s, "size": 12, "bold": True, "color": BLUE, "align": PP_ALIGN.CENTER}])
    text(slide, M + 10, y + 39, lw_ - 20, 26, [
        {"t": "코드부터 쓰지 않았습니다. 문제와 업무규칙을 먼저 정의하고 시험 기준을 확정한 뒤, "
              "구현·디버깅·시험을 AI와 반복했습니다. 실사용 결함은 보고 당일 수정했습니다.",
         "size": 13, "color": MUTED, "after": 2.5, "spacing": 1.2},
        {"t": f"실사용 중 접수한 개선 요청 {DATA['todo_all']} 중 {DATA['todo_done']} 반영, "
              f"{DATA['todo_wait']}은 판단 대기.",
         "size": 13, "color": MUTED, "spacing": 1.2},
    ])

    # 오른쪽 : 향후 계획
    rx = M + lw_ + 10
    rect(slide, rx, y, rw_, LH, fill=NAVY, radius=0.08)
    text(slide, rx + 10, y + 8, rw_ - 20, 8,
         [{"t": "앞으로", "size": 15, "bold": True, "color": SKY}])
    plans = [
        ("1단계", "사내 서버로 확장 — 여러 팀장 공동 사용, 사용자·과제번호·권한 체계 표준화"),
        ("2단계", "주간회의 선정 과제와 의사결정 이력까지 관리 — 의사결정 지원 시스템으로"),
        ("3단계", "축적된 데이터로 업무현황·이슈 분석, 반복 보고 업무를 AI가 지원"),
        ("확산", "이 6단계 방법론을 표준화해 다른 현업도 직접 업무도구를 만들도록"),
    ]
    py = y + 18
    for i, (k, v) in enumerate(plans):
        rect(slide, rx + 10, py + 0.3, 15, 8, fill=BLUE, radius=0.4)
        text(slide, rx + 10, py + 1.9, 15, 7,
             [{"t": k, "size": 11, "bold": True, "color": WHITE, "align": PP_ALIGN.CENTER}])
        text(slide, rx + 30, py + 0.4, rw_ - 42, 13,
             [{"t": v, "size": 12, "color": WHITE, "spacing": 1.12}])
        py += 11.6
    y += LH + 6

    # ── 꼬리말 ────────────────────────────────────────────────────────────
    rect(slide, M, y, W, 0.6, fill=LINE)
    text(slide, M, y + 4, W, 16, [
        {"t": "규모·기간·시험 건수는 2026-09-04 실측치입니다. 외주 비용은 기능점수 기반 추정이며, "
              "단가(시장 단가·직접인건비·FP 단가)는 가정값이므로 「소프트웨어사업 대가산정 가이드」와 "
              "「SW기술자 평균임금」 최신 공표치로 대체해 확인하시기 바랍니다.",
         "size": 12, "color": FAINT, "spacing": 1.25},
    ])

    print(f"마지막 요소 아래 여백 {PH - y - 22:.1f}mm")
    path = OUT / "AX성과공유회_포스터_개발을기다리지않았다.pptx"
    prs.save(path)
    return path


if __name__ == "__main__":
    p = build()
    print(p, f"bottom≈{PH}mm")
